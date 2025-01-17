"""
PDF to PPTX Converter

Author: Jeongkyu Shin <inureyes@gmail.com>

A Python script to convert PDF files to PowerPoint presentations, with each page of the PDF being an image in a slide.
The script supports conversion to both JPG and PNG image formats.

Features:
- Converts each page of a PDF file to an image and places it in a PowerPoint slide.
- Supports JPG and PNG image formats.
- Automatically handles file naming and temporary folders.
- Displays progress bars for conversion and slide creation.
- Automatically cleans up temporary files after creating the presentation.
- On macOS, if Microsoft PowerPoint is installed, you can directly input a PPTX file to flatten it.

License:
This project is licensed under the MIT License. See the LICENSE file for details.

Acknowledgements:
This project uses the following libraries:
- pdf2image
- python-pptx
- tqdm
- Poetry

Issues and Contributions:
Feel free to submit issues and contribute to the project on GitHub.
"""

import os
import sys
import shutil
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm
import argparse
import subprocess

def is_powerpoint_installed():
    try:
        result = subprocess.run(["osascript", "-e", 'id of application "Microsoft PowerPoint"'], capture_output=True, text=True)
        if "com.microsoft.Powerpoint" in result.stdout:
            return True
        else:
            return False
    except subprocess.CalledProcessError:
        return False

def convert_ppt_to_pdf_mac(input_ppt_path, output_pdf_path):
    input_ppt_path = os.path.abspath(input_ppt_path)
    output_pdf_path = os.path.abspath(output_pdf_path)
    script = f'''
    tell application "Microsoft PowerPoint"
        open POSIX file "{input_ppt_path}"
        set theDoc to active presentation
        save theDoc in POSIX file "{output_pdf_path}" as save as PDF
        close theDoc saving no
    end tell
    '''
    process = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
    if process.returncode != 0:
        print("AppleScript Error:", process.stderr)
        sys.exit(1)

def get_output_pptx_path(base_name):
    i = 1
    output_pptx_path = f"{base_name} (flatten).pptx"
    while os.path.exists(output_pptx_path):
        output_pptx_path = f"{base_name} (flatten)_{i}.pptx"
        i += 1
    return output_pptx_path

def main():
    parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint with images")
    parser.add_argument("input_path", help="Path to the input PPT or PDF file")
    parser.add_argument("--format", choices=["png", "jpg"], default="jpg", help="Image format for conversion (default: jpg)")
    args = parser.parse_args()

    input_path = args.input_path
    image_format = args.format

    if not os.path.exists(input_path):
        print(f"Error: File '{input_path}' does not exist.")
        sys.exit(1)

    base_name, ext = os.path.splitext(os.path.basename(input_path))
    created_pdf = False

    try:
        if ext.lower() in [".ppt", ".pptx"]:
            if not is_powerpoint_installed():
                print("Error: Microsoft PowerPoint is not installed on this Mac.")
                sys.exit(1)
            input_pdf_path = os.path.join(os.path.dirname(input_path), f"{base_name}.pdf")
            print("Converting PowerPoint to PDF...")
            convert_ppt_to_pdf_mac(input_path, input_pdf_path)
            created_pdf = True
        elif ext.lower() == ".pdf":
            input_pdf_path = input_path
        else:
            print("Error: Unsupported file format. Please provide a PPT, PPTX, or PDF file.")
            sys.exit(1)

        image_folder_path = f"output_images_{base_name}"

        if not os.path.exists(image_folder_path):
            print("Creating temporary folder for images...")
            os.makedirs(image_folder_path)

        # Convert PDF pages to images. DPI can be adjusted for quality and final pptx file size.
        images = convert_from_path(input_pdf_path, dpi=300)

        print(f"Converting PDF pages to {image_format} images...")
        num_pages = len(images)
        num_digits = len(str(num_pages))
        for i, image in enumerate(tqdm(images, desc="Pages", unit="page", ncols=80, bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')):
            image_path = os.path.join(image_folder_path, f"page_{str(i + 1).zfill(num_digits)}.{image_format}")
            if image_format == 'jpg':
                image_format_upper = 'JPEG'
            else:
                image_format_upper = image_format.upper()
            image.save(image_path, image_format_upper)

        # Create a PowerPoint presentation and add slides with images. Be careful with the slide dimensions.
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        image_files = sorted([f for f in os.listdir(image_folder_path) if f.endswith(f'.{image_format}')])
        print("Adding images to the PowerPoint presentation...")
        for image_filename in tqdm(image_files, desc="Slides", unit="slide", ncols=80, bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}'):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            img_path = os.path.join(image_folder_path, image_filename)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        output_pptx_path = get_output_pptx_path(base_name)
        prs.save(output_pptx_path)
        print(f"New presentation saved as {output_pptx_path}")

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
    finally:
        if os.path.exists(image_folder_path):
            print("Deleting temporary folder for images...")
            shutil.rmtree(image_folder_path)

        if created_pdf and os.path.exists(input_pdf_path):
            print(f"Deleting temporary PDF file {input_pdf_path}...")
            os.remove(input_pdf_path)

if __name__ == "__main__":
    main()